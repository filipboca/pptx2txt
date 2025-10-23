import os
import re
from pptx import Presentation

# ---------- TEXT CLEANUP HELPERS ----------

def normalize_newlines(text):
    """
    Convert all newline-like Unicode characters to '\n'.
    """
    return (text
        .replace('\r\n', '\n')
        .replace('\r', '\n')
        .replace('\u2028', '\n')
        .replace('\u2029', '\n')
        .replace('\v', '\n')
        .replace('\f', '\n')
    )

def clean_text(text):
    """
    Normalize and clean PowerPoint text content:
    - Normalize newlines
    - Remove invisible or unwanted characters
    - Collapse excessive whitespace
    """
    text = normalize_newlines(text)

    # Remove zero-width spaces, BOM, non-breaking spaces, soft hyphens
    text = re.sub(r'[\u200B-\u200D\uFEFF]', '', text)   # zero-widths
    text = text.replace('\u00A0', ' ')                  # non-breaking space → normal space
    text = text.replace('\u00AD', '')                   # soft hyphen

    # Remove control characters except tab (0x09) and newline (0x0A)
    text = re.sub(r'[^\x09\x0A\x20-\x7E\u00A0-\u024F]', '', text)

    # Collapse multiple spaces (but preserve single \n)
    text = re.sub(r' {2,}', ' ', text)

    # Normalize too many blank lines
    text = re.sub(r'\n{3,}', '\n\n', text)

    return text.strip()


# ---------- MAIN EXTRACTION LOGIC ----------

def extract_text_from_pptx(path):
    """
    Extract all text (with formatting cleaned) from a .pptx file.
    """
    prs = Presentation(path)
    all_text = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                cleaned = clean_text(shape.text)
                if cleaned:
                    slide_text.append(cleaned)

        if slide_text:
            slide_block = f"\n".join(slide_text)
            all_text.append(slide_block)

    return "\n\n".join(all_text)


def convert_pptx_folder(input_folder, output_folder=None):
    """
    Process all .pptx files in a folder (recursively) and write .txt outputs.
    """
    if output_folder is None:
        output_folder = input_folder

    os.makedirs(output_folder, exist_ok=True)
    pptx_files = []

    for root, _, files in os.walk(input_folder):
        for file in files:
            if file.lower().endswith(".pptx"):
                pptx_files.append(os.path.join(root, file))

    if not pptx_files:
        print("No .pptx files found in the folder.")
        return

    for pptx_path in pptx_files:
        relative_name = os.path.relpath(pptx_path, input_folder)
        txt_name = os.path.splitext(relative_name)[0] + ".txt"
        output_path = os.path.join(output_folder, txt_name)

        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        print(f"📄 Processing: {relative_name} ...")
        try:
            text = extract_text_from_pptx(pptx_path)
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(text)
            print(f"Saved to: {output_path}")
        except Exception as e:
            print(f"Error processing {pptx_path}: {e}")


# ---------- ENTRY POINT ----------

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python pptx_to_text_batch.py <input_folder> [output_folder]")
        sys.exit(1)

    input_folder = sys.argv[1]
    output_folder = sys.argv[2] if len(sys.argv) > 2 else None

    convert_pptx_folder(input_folder, output_folder)
