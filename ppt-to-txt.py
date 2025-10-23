from pptx import Presentation
import sys
import os
import re

def normalize_newlines(text):
    """
    Convert all newline variants to '\n' for consistency.
    """
    return (text
        .replace('\r\n', '\n')  # Windows
        .replace('\r', '\n')    # Old Mac
        .replace('\u2028', '\n')  # Unicode line separator
        .replace('\u2029', '\n')  # Unicode paragraph separator
        .replace('\v', '\n')      # Vertical tab
        .replace('\f', '\n')      # Form feed
    )

def clean_text(text):
    """
    Normalize newlines and remove unwanted Unicode control characters,
    extra spaces, and invisible formatting marks.
    """
    text = normalize_newlines(text)

    # Remove zero-width and non-printing Unicode characters
    text = re.sub(r'[\u200B\u200C\u200D\uFEFF]', '', text)  # zero-widths
    text = re.sub(r'[\u00A0]', ' ', text)                   # non-breaking space → normal space
    text = re.sub(r'[\u00AD]', '', text)                    # soft hyphen
    text = re.sub(r'[^\x09\x0A\x20-\x7E\u00A0-\u024F]', '', text)  # remove control chars except tab/newline

    # Replace multiple spaces with a single space (but preserve newlines)
    text = re.sub(r'[ ]{2,}', ' ', text)

    # Normalize double newlines (avoid accidental blank gaps)
    text = re.sub(r'\n{3,}', '\n\n', text)

    return text.strip()

def extract_text_from_pptx(path):
    prs = Presentation(path)
    all_text = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                #slide_text.append(normalize_newlines(shape.text))
                cleaned = clean_text(shape.text)
                if cleaned:
                    slide_text.append(cleaned)

        if slide_text:
            # Join text boxes with newlines, and mark slide boundaries
            all_text.append(f"\n".join(slide_text))

    return "\n\n".join(all_text)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx_text.py <input.pptx> [output.txt]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        sys.exit(1)

    if len(sys.argv) >= 3:
        output_path = sys.argv[2]
    else:
        output_path = os.path.splitext(pptx_path)[0] + ".txt"

    text = extract_text_from_pptx(pptx_path)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)

    print(f"✅ Text extracted and saved to {output_path}")
